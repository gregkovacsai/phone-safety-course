#!/usr/bin/env python3
"""Build a polished PowerPoint from Elli's Phone Safety Course."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Colors ──
BG = RGBColor(0x1a, 0x1a, 0x2e)
BG_SURFACE = RGBColor(0x16, 0x21, 0x3e)
WHITE = RGBColor(0xff, 0xff, 0xff)
DIM = RGBColor(0xa7, 0xa9, 0xbe)
PURPLE = RGBColor(0xa8, 0x55, 0xf7)
TEAL = RGBColor(0x06, 0xb6, 0xd4)
PINK = RGBColor(0xe0, 0x56, 0xa0)
WARN = RGBColor(0xff, 0x6b, 0x6b)
GREEN = RGBColor(0x2c, 0xb6, 0x7d)

FONT = 'Calibri'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

slide_number_counter = [0]

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    slide_number_counter[0] += 1
    # Slide number
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_number_counter[0])
    p.font.size = Pt(10)
    p.font.color.rgb = DIM
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT
    return slide

def add_text(slide, left, top, width, height, text, font_size=28, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_rich_text(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT):
    """runs = list of (text, size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, (text, size, color, bold) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
    return tf

def add_bullets(slide, left, top, width, height, items, font_size=26, icon_color=PURPLE):
    """items = list of (emoji, text)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (emoji, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{emoji}  {text}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE
        run.font.name = FONT
    return tf

def add_callout(slide, left, top, width, height, label, text, label_color=PINK, border_color=PURPLE):
    """Callout box with colored left border effect."""
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_SURFACE
    shape.line.fill.background()
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Emu(Inches(0.05).emu), Inches(0.08), height - Emu(Inches(0.1).emu))
    bar.fill.solid()
    bar.fill.fore_color.rgb = border_color
    bar.line.fill.background()
    # Label
    add_text(slide, left + Inches(0.25), top + Inches(0.1), width - Inches(0.4), Inches(0.3),
             label.upper(), font_size=12, color=label_color, bold=True)
    # Text
    add_text(slide, left + Inches(0.25), top + Inches(0.38), width - Inches(0.4), height - Inches(0.5),
             text, font_size=18, color=DIM)

def add_discussion(slide, left, top, width, height, question):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1f, 0x1a, 0x3e)
    shape.line.color.rgb = PURPLE
    shape.line.width = Pt(2)
    add_text(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.4),
             "💬 DISCUSSION PROMPT", font_size=19, color=PURPLE, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.55), width - Inches(0.6), height - Inches(0.7),
             question, font_size=28, color=WHITE)

def title_slide(emoji, title, subtitle, module_label=None):
    slide = add_slide()
    y = Inches(1.0)
    if module_label:
        add_text(slide, Inches(1), y, Inches(11.333), Inches(0.4),
                 module_label.upper(), font_size=14, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
        y += Inches(0.5)
    add_text(slide, Inches(1), y, Inches(11.333), Inches(1.0),
             emoji, font_size=72, alignment=PP_ALIGN.CENTER)
    y += Inches(1.1)
    add_text(slide, Inches(1), y, Inches(11.333), Inches(1.0),
             title, font_size=48, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
    y += Inches(1.0)
    add_text(slide, Inches(1.5), y, Inches(10.333), Inches(1.5),
             subtitle, font_size=26, color=DIM, alignment=PP_ALIGN.CENTER)
    return slide

def content_slide(emoji, title, body_items=None, body_text=None, callouts=None, note=None):
    slide = add_slide()
    y = Inches(0.4)
    if emoji:
        add_text(slide, Inches(1), y, Inches(11.333), Inches(0.8),
                 emoji, font_size=52, alignment=PP_ALIGN.CENTER)
        y += Inches(0.8)
    add_text(slide, Inches(0.8), y, Inches(11.733), Inches(0.6),
             title, font_size=36, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
    y += Inches(0.7)
    if body_text:
        add_text(slide, Inches(1.5), y, Inches(10.333), Inches(0.8),
                 body_text, font_size=28, color=DIM, alignment=PP_ALIGN.CENTER)
        y += Inches(0.8)
    if body_items:
        add_bullets(slide, Inches(1.5), y, Inches(10.333), Inches(4.5), body_items)
        y += Inches(len(body_items) * 0.38)
    if callouts:
        for ctype, label, text in callouts:
            lc = PINK if ctype == 'rt' else (WARN if ctype == 'warn' else GREEN)
            bc = PURPLE if ctype == 'rt' else (WARN if ctype == 'warn' else GREEN)
            add_callout(slide, Inches(1.5), y, Inches(10.333), Inches(1.0), label, text, lc, bc)
            y += Inches(1.1)
    if note:
        add_text(slide, Inches(1.5), y, Inches(10.333), Inches(0.5),
                 note, font_size=18, color=PURPLE, alignment=PP_ALIGN.CENTER)
    return slide

def discussion_slide(question):
    slide = add_slide()
    add_discussion(slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5), question)
    return slide

def quiz_slide(questions, slide_title="📝 Quiz"):
    """questions = list of (question, options, correct_idx, explanation)"""
    slide = add_slide()
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11.733), Inches(0.6),
             slide_title, font_size=34, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
    y = Inches(1.0)
    for q, opts, ci, expl in questions:
        add_text(slide, Inches(1.2), y, Inches(10.9), Inches(0.5),
                 q, font_size=26, color=WHITE, bold=True)
        y += Inches(0.45)
        for i, o in enumerate(opts):
            marker = "✅ " if i == ci else "○ "
            c = GREEN if i == ci else DIM
            add_text(slide, Inches(1.6), y, Inches(10.5), Inches(0.3),
                     marker + o, font_size=16, color=c)
            y += Inches(0.28)
        # Speaker notes for explanation
        y += Inches(0.15)
    # Add explanations to speaker notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    for q, opts, ci, expl in questions:
        p = notes_tf.add_paragraph()
        p.text = f"Q: {q}\nA: {expl}\n"
    return slide

# ═══════════════════════════════════════════
# BUILD ALL SLIDES
# ═══════════════════════════════════════════

# 1 - Welcome
title_slide("📱", "Welcome, Elli!", "Your Personal Guide to Being Smart, Safe & In Control\n\nMade just for you, with love from Mom & Dad ❤️")

# 2 - First Things First
content_slide("🎉", "First Things First, Elli",
    body_text="Getting your own phone is a big deal. This course isn't a punishment — it's your preparation.\n\nThink of it like getting a driver's license. You wouldn't just hop in a car without learning the rules, right?",
    callouts=None)

# 3 - We Trust You
content_slide("💪", "We Trust You, Elli",
    body_text="The fact that you're getting a phone means your parents believe you're ready for more responsibility.\n\nThis course is about giving you the knowledge and tools to handle it like a boss.")

# 4 - What a Phone Really Means
content_slide("⚖️", "What a Phone Really Means", [
    ("🌍", "Access to the entire world — the good AND the bad"),
    ("🗣️", "A direct line to anyone, anywhere, anytime"),
    ("📸", "A camera that can capture (and share) anything"),
    ("🧠", "A tool that can make your life better — or worse"),
    ("💡", "It all depends on how YOU use it"),
])

# 5 - What You'll Learn
content_slide("🎯", "What You'll Learn Today", [
    ("🔧", "Building healthy tech habits"),
    ("🛡️", "Spotting scams and staying safe from predators"),
    ("🔒", "Protecting your personal information"),
    ("🧠", "Protecting your mental health"),
    ("🤖", "Navigating AI smartly"),
    ("📊", "Outsmarting the algorithm"),
    ("📱", "Social media — the real talk"),
    ("🌟", "Using your phone to level up your life"),
])

# 6 - Discussion
discussion_slide("Before we start: Elli, what are you MOST excited about having your own phone? What (if anything) makes you nervous about it?")

# ── MODULE 2: TECH HABITS ──
title_slide("⏰", "Building Tech Habits That Stick",
    "You control the phone.\nThe phone does NOT control you.",
    "Module 2 — Tech Habits")

# Screen Time
content_slide("📊", "Screen Time Awareness",
    body_text="The average teen spends 7+ hours per day on screens (outside of school).\n\nThat's almost a full-time job! Being aware of your usage is the first step to staying in control.")

# Set Boundaries
content_slide("⏱️", "Set Your Own Boundaries", [
    ("📱", "Use Screen Time / Digital Wellbeing to track your usage"),
    ("⏰", "Set daily app limits (especially TikTok and YouTube — they're time machines!)"),
    ("🎯", 'Ask yourself: "Am I using my phone on purpose, or just because I\'m bored?"'),
    ("📝", "Try a \"screen time journal\" for the first week — you'll be surprised!"),
])

# Notifications
content_slide("🔔", "Tame Your Notifications",
    body_text="Every buzz, ding, and banner is designed to pull you back in.",
    body_items=[
        ("🔇", "Turn off notifications for most apps"),
        ("⭐", "Only keep notifications for calls, texts from family, and essentials"),
        ("🧘", 'Try "Do Not Disturb" mode during homework and meals'),
        ("💡", "YOU decide when to check your phone — not the app"),
    ])

# Sleep
content_slide("😴", "Sleep Hygiene",
    body_text="Blue light and late-night scrolling wreck your sleep — and sleep is everything at your age.",
    body_items=[
        ("🚫", "No phone in your bedroom at night"),
        ("🔌", "Phone charges in the kitchen/living room overnight"),
        ("⏰", "Screen curfew: 1 hour before bed"),
        ("😌", "Your brain needs downtime to process the day"),
    ],
    callouts=[("rt", "Real Talk", "Studies show teens who keep phones in their bedroom get 30 minutes LESS sleep per night on average. Over a year, that's 180 hours of lost sleep.")])

# Phone Stack
content_slide("🃏", "The Phone Stack Game",
    body_text="When hanging out with friends or family:",
    body_items=[
        ("📱", "Everyone stacks their phones in the middle of the table"),
        ("🚫", "First person to grab their phone loses (pays for dessert, does a dare, etc.)"),
        ("🤝", "It's about being PRESENT with the people in front of you"),
    ],
    note="The people in front of you always matter more than the people on your screen.")

# Phone-Free Zones
content_slide("📵", "Phone-Free Zones & Times", [
    ("🍽️", "Meals — always phone-free"),
    ("📚", "Homework time — phone in another room"),
    ("🛏️", "Bedroom at night"),
    ("⛪", "Family events, ceremonies, gatherings"),
    ("🚗", "Car rides (try talking instead!)"),
], callouts=[("tip", "Pro Tip", "Making these habits now means they'll feel natural forever. It's way harder to break bad habits than to build good ones from the start.")])

# Discussion
discussion_slide("Elli, what phone-free zones and times make sense for your family? Discuss and agree on at least 3 together.")

# ── MODULE 3: SAFETY ──
title_slide("🛡️", "Watching Out for Scams & Predators",
    "This is the most important module.\nPlease take it seriously, Elli.",
    "Module 3 — Safety")

# Scams
content_slide("🎣", "How Scams Work",
    body_text="Scammers use psychological tricks to get you to act without thinking:",
    body_items=[
        ("⏰", 'Urgency — "Act NOW or lose your account!"'),
        ("🎁", 'Too good to be true — "You won a free iPhone!"'),
        ("😨", 'Fear — "Your account has been hacked!"'),
        ("❤️", 'Emotion — "Help this sick puppy!"'),
        ("🎭", "Impersonation — pretending to be a friend, company, or authority"),
    ])

# Scams Targeting Teens
s = add_slide()
add_text(s, Inches(0.8), Inches(0.3), Inches(11.733), Inches(0.6),
         "🎯 Scams Targeting Teens", font_size=36, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_callout(s, Inches(1.2), Inches(1.1), Inches(10.9), Inches(1.3),
    "Fake Giveaway",
    'You see a TikTok post: "🎉 FREE iPhone! Just follow, like, share, and enter your email + address!" — This is ALWAYS a scam.',
    WARN, WARN)
add_callout(s, Inches(1.2), Inches(2.6), Inches(10.9), Inches(1.3),
    "Phishing Text",
    '"Your Snapchat account will be deleted in 24 hours. Click here to verify." — Fake link. Real Snapchat would never text you this.',
    WARN, WARN)
add_callout(s, Inches(1.2), Inches(4.1), Inches(10.9), Inches(1.3),
    "Cash App Flip",
    '"Send me $50 and I\'ll flip it to $500 with this money hack!" — Nobody can magically multiply money. This is theft.',
    WARN, WARN)

# Predators
content_slide("⚠️", "Online Predators",
    body_text="Predators are adults who try to build relationships with young people online to exploit them. They are very good at what they do.\n\nThey don't look like villains. They seem friendly, understanding, and cool.")

# Grooming
content_slide("🎭", "How Grooming Works",
    body_text="Grooming is a step-by-step process:",
    body_items=[
        ("1️⃣", "Targeting — They find someone who seems lonely, insecure, or seeking attention"),
        ("2️⃣", 'Building trust — "You\'re so mature for your age" / "I totally get you"'),
        ("3️⃣", "Filling a need — Compliments, gifts, attention, \"understanding\""),
        ("4️⃣", 'Isolating — "Don\'t tell your parents, they wouldn\'t understand"'),
        ("5️⃣", "Desensitizing — Gradually introducing inappropriate topics"),
        ("6️⃣", "Exploiting — Asking for photos, meetups, or favors"),
    ])

# Red Flags
content_slide("🚩", "Red Flags in Online Conversations", [
    ("🚩", "They ask you to keep the friendship a secret"),
    ("🚩", 'They say "you\'re so mature for your age"'),
    ("🚩", "They ask personal questions quickly (where you live, what school)"),
    ("🚩", "They send you gifts or offer money"),
    ("🚩", "They try to move conversation to a private platform"),
    ("🚩", "They ask for photos (especially selfies)"),
    ("🚩", "They get upset or guilt-trip you if you say no"),
    ("🚩", 'They claim to be a teen but something feels "off"'),
])

# Never Share
content_slide("🔐", "NEVER Share These Online", [
    ("🏠", "Home address"),
    ("🏫", "School name"),
    ("📍", "Current location"),
    ("📸", "Personal photos"),
    ("📞", "Phone number"),
    ("👤", "Full name"),
    ("🎂", "Birthday"),
    ("🔑", "Passwords"),
], note="Not with strangers. Not with online friends. Not in games. Not ever.")

# If Something Feels Wrong
s = add_slide()
add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         "🆘 If Something Feels Wrong...", font_size=36, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x15, 0x2a, 0x1f)
shape.line.color.rgb = GREEN
shape.line.width = Pt(3)
add_text(s, Inches(2), Inches(1.7), Inches(9.333), Inches(0.4),
         "THE #1 RULE", font_size=14, color=GREEN, bold=True)
add_text(s, Inches(2), Inches(2.2), Inches(9.333), Inches(0.6),
         "Tell a parent IMMEDIATELY.", font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(3.0), Inches(9.333), Inches(2.5),
         "You will NOT be in trouble. Not now, not ever.\nEven if you think you did something wrong.\nEven if someone told you not to tell.\nEven if you're embarrassed.\n\nElli, Mom and Dad are ALWAYS on your team. 💚",
         font_size=26, color=DIM, alignment=PP_ALIGN.CENTER)

# Real Talk Scenario
s = add_slide()
add_callout(s, Inches(1.2), Inches(1.0), Inches(10.9), Inches(2.5),
    "Real Talk Scenario",
    "You're playing an online game and someone who says they're 14 starts chatting with you. They're really nice and funny. After a week, they ask what school you go to and if you want to video call — but ask you not to tell your parents because \"they might not let us be friends.\"",
    PINK, PURPLE)
add_text(s, Inches(1.5), Inches(4.0), Inches(10.333), Inches(1.0),
         "Every single thing about this is a red flag.\nA real friend would never ask you to hide the friendship.",
         font_size=28, color=WARN, bold=True, alignment=PP_ALIGN.CENTER)

# ── EXPLICIT CONTENT SECTION ──
content_slide("🔞", "Explicit Content Online",
    body_text="Elli, this is an awkward topic but an important one. Pornography and explicit content exist all over the internet — and it can show up even when you're not looking for it.\n\nPop-up ads, links in group chats, search results, social media — it can appear unexpectedly on almost any platform.")

content_slide("🚫", "Why This Content Is Harmful", [
    ("🧠", "It creates unrealistic and harmful expectations about relationships and bodies"),
    ("💔", "It can warp your understanding of what healthy relationships look like"),
    ("😰", "It can make you feel confused, uncomfortable, or anxious"),
    ("📱", "It's designed for adults and does NOT reflect real life — not even close"),
    ("🌱", "Your brain is still developing, and this content can genuinely affect how you see yourself and others"),
])

content_slide("🛑", "If You Accidentally See Something", [
    ("❌", "Close it immediately — don't keep looking out of curiosity"),
    ("🗣️", "Tell Mom or Dad — you will NOT be in trouble, we promise"),
    ("💚", "No shame — accidentally seeing something doesn't mean you did anything wrong"),
    ("🚫", "Don't share it — sharing explicit content involving minors is actually illegal"),
    ("🔍", "Don't go looking for more — curiosity is normal, but this content is genuinely harmful"),
], callouts=[("green", "Remember", "Seeing something explicit by accident is NOT your fault. It happens to almost everyone online. What matters is what you do next.")])

content_slide("⚠️", "Where Explicit Content Can Appear", [
    ("💬", "Group chats and DMs (people send links or images)"),
    ("🔍", "Search engines and image searches (even innocent searches)"),
    ("📱", "Social media feeds — TikTok, YouTube, even Pinterest"),
    ("🎮", "Online games and gaming chat platforms"),
    ("📺", "Streaming sites with inadequate age filters"),
    ("💻", "Pop-up ads on free websites"),
], note="Being aware of where it can appear helps you be prepared if it does.")

content_slide("🛡️", "Parental Controls Are on Your Side",
    body_text="Elli, the parental controls on your devices are there to help you, not spy on you.",
    body_items=[
        ("🔒", "They filter out content that no kid should have to see"),
        ("💚", "They're like a seatbelt — a safety tool, not a punishment"),
        ("🤝", "As you get older and show responsibility, settings can be adjusted together"),
        ("🗣️", "If a control blocks something you need for school, just ask"),
    ],
    callouts=[("tip", "The Big Picture", "We're not trying to control your every move. We're trying to make sure you don't stumble into stuff that could genuinely hurt you. That's our job as parents. 💚")])

discussion_slide("Elli, have you ever encountered anything online that made you uncomfortable? What would you do if an online stranger asked you to keep a secret from us? And remember — if you ever see something explicit, you can always come to us without any judgment.")

# ── MODULE 4: PRIVACY ──
title_slide("🔒", "Protecting Your Personal Information",
    "Your data is valuable.\nTreat it like treasure, Elli.",
    "Module 4 — Privacy")

content_slide("💎", "What Counts as Personal Info?", [
    ("👤", "Full name, birthday, age"),
    ("🏠", "Address, phone number, email"),
    ("🏫", "School name, team/club names"),
    ("📸", "Photos of yourself, your home, your ID"),
    ("📍", "Location data (where you are right now)"),
    ("🔑", "Passwords, security questions"),
    ("💳", "Financial info (even your parents')"),
    ("🧬", "Health info, family details"),
], note="If it can be used to find you, identify you, or impersonate you — protect it.")

content_slide("⚙️", "Privacy Settings Matter", [
    ("🔒", "Set ALL social media accounts to private"),
    ("📍", "Turn off location sharing in apps (except with family)"),
    ("📷", "Disable camera/microphone access for apps that don't need it"),
    ("🔍", "Google yourself — see what's already out there"),
    ("📋", "Review app permissions regularly"),
], callouts=[("tip", "Pro Tip", 'When installing a new app, ask: "Does this app REALLY need access to my contacts and location?" If the permissions don\'t make sense, don\'t install it.')])

content_slide("🏷️", '"Free" Apps Aren\'t Free',
    body_text="If you're not paying for the product...\n\nYOU are the product.\n\nFree apps like TikTok, Pinterest, and YouTube make money by collecting your data and selling it to advertisers. Every tap, scroll, search, and like is tracked and monetized.")

content_slide("🔑", "Password Hygiene", [
    ("✅", "Use long, unique passwords for every account (12+ characters)"),
    ("✅", "Mix uppercase, lowercase, numbers, and symbols"),
    ("✅", "Use a password manager (ask Mom/Dad to help set one up)"),
    ("❌", "NEVER use your name, birthday, or pet's name"),
    ("❌", "NEVER reuse passwords across sites"),
    ("❌", "NEVER share passwords with friends (even best friends)"),
], callouts=[("tip", "Easy Password Trick", 'Think of a sentence: "My cat Luna loves 3 treats at bedtime!" → McLl3t@b! — long, random, and easy for YOU to remember.')])

content_slide("👣", "Your Digital Footprint",
    body_text="Everything you post, like, comment, search for, or share online creates a permanent trail.",
    body_items=[
        ("🏫", "College admissions officers check social media"),
        ("💼", "Future employers will Google you"),
        ("📸", "Screenshots mean deleted posts aren't really deleted"),
        ("⏰", "Something you post at 13 can follow you at 30"),
    ],
    note="The internet never forgets, Elli.")

discussion_slide("What's one thing you didn't realize counted as 'personal information' before today? How will you decide which apps to give permissions to?")

# ── MODULE 5: MENTAL HEALTH ──
title_slide("🧠", "Protecting Your Mental Health",
    "Your phone should make your life better.\nIf it's not, something needs to change.",
    "Module 5 — Mental Health")

content_slide("🪞", "The Comparison Trap",
    body_text="What you see on social media is a highlight reel, not real life.",
    body_items=[
        ("📸", "People post their best moments, not their bad days"),
        ("✨", "Filters, editing, and staging make everything look perfect"),
        ("🎭", "Even influencers have acne, bad hair days, and insecurities"),
        ("💔", "Comparing your behind-the-scenes to someone else's highlight reel is unfair to YOU"),
    ])

content_slide("📌", "A Note About Pinterest",
    body_text="Pinterest is generally a great app for inspiration and creativity, Elli. But there's a flip side:",
    body_items=[
        ("🪞", 'Boards full of "perfect" bodies, rooms, outfits can create unrealistic standards'),
        ("📊", "The more you save certain types of content, the more it shows you — this can spiral"),
        ("💚", "Curate intentionally — fill your boards with things that inspire YOU"),
        ("🎨", "Use it for creative projects, book ideas, art inspiration — that's where Pinterest shines!"),
    ],
    callouts=[("tip", "Healthy Pinterest Habit", "If you notice you're feeling bad about yourself after scrolling Pinterest, that's a sign to curate your boards differently. Pin content that sparks JOY and creativity.")])

content_slide("🪞", "AI Face Filters — The Fun Lie",
    body_text="Filters that change your face are everywhere. Fun — but they can mess with how you see yourself.",
    body_items=[
        ("📸", "TikTok, Snapchat, and Instagram filters change how you look in real-time"),
        ("🧠", 'Using them constantly can make your REAL face feel "wrong" — called "Snapchat dysmorphia"'),
        ("🪞", "You start comparing yourself to a filtered version of YOU that doesn't exist"),
        ("💔", 'Some filters make you look older, thinner, or more "perfect" — that\'s not enhancement, it\'s erasure'),
        ("😄", "Silly/fun filters (dog ears, rainbow vomit)? Go for it!"),
        ("⚠️", "Filters that change your actual features? Be careful how often you use them"),
        ("❤️", "If you feel \"worse\" without a filter — that's the filter's damage, not your face"),
    ],
    note='"Your real face is the one people love, Elli."')

content_slide("😰", "FOMO — Fear of Missing Out",
    body_text="That sinking feeling when everyone seems to be having fun without you.",
    body_items=[
        ("📱", "Seeing friends hang out without you on Snapchat or TikTok"),
        ("🎉", "Feeling like everyone's life is more exciting"),
        ("😔", "Checking your phone constantly for updates"),
    ],
    callouts=[("green", "The Truth", "FOMO is manufactured by social media. People only post when things look fun. JOMO (Joy of Missing Out) is a real thing — being happy with what YOU'RE doing right now.")])

content_slide("😢", "Cyberbullying",
    body_text="Bullying that happens through phones and online:",
    body_items=[
        ("💬", "Mean messages, comments, or DMs"),
        ("📸", "Sharing embarrassing photos without permission"),
        ("🚫", "Deliberately excluding someone in group chats"),
        ("👤", "Creating fake accounts to harass someone"),
        ("📢", "Spreading rumors online"),
        ("🔄", "Screenshotting private conversations to embarrass someone"),
    ])

content_slide("🛡️", "If You're Being Cyberbullied", [
    ("📸", "Screenshot everything (evidence matters)"),
    ("🚫", "Don't respond (they want a reaction)"),
    ("🔒", "Block the person"),
    ("🗣️", "Tell a parent or trusted adult"),
    ("📝", "Report it on the platform"),
], callouts=[("green", "If You See It Happening to Someone Else...", "Don't be a bystander. Don't join in or share it. Stand up for them, or tell an adult. Being kind online takes courage — and it matters more than you know.")])

content_slide("👥", "Group Chat Pressure",
    body_text="Group chats can be fun — but they can also get weird fast.",
    body_items=[
        ("🤐", "Pressure to agree with the group or pile on someone"),
        ("📸", "Someone shares something inappropriate — now everyone's seen it"),
        ("🚪", "Being added to chats without permission"),
        ("😬", "Drama that spirals fast because everyone's watching"),
        ("🗣️", "People say things in groups they'd never say 1-on-1"),
        ("🚪", "You can LEAVE any group chat that makes you uncomfortable"),
        ("🤫", "You don't have to respond to everything"),
        ("📸", "If something crosses a line, screenshot and tell a parent"),
        ("🔇", "Mute chats that stress you out"),
    ])

content_slide("🔄", "The Dopamine Loop",
    body_text="Here's why you can't stop scrolling:",
    body_items=[
        ("🧪", 'Every like, comment, and new post gives your brain a tiny hit of dopamine (the "feel good" chemical)'),
        ("🎰", "Apps are designed like slot machines — you keep scrolling hoping for the next reward"),
        ("📱", "This is NOT an accident. Billions of dollars are spent making apps as addictive as possible"),
        ("🧠", "Knowing this gives you power over it"),
    ])

content_slide("📉", "Real Stats on Teens & Phones", [
    ("📊", "Teens who spend 5+ hrs/day on social media are 3x more likely to report depression"),
    ("😴", "70% of teens say social media makes them feel worse about their appearance"),
    ("📱", "1 in 3 teens say they wish they could go back to life before social media"),
    ("💚", "BUT — teens who use phones for connection and creativity report higher well-being"),
], note="It's not about having a phone. It's about HOW you use it.")

content_slide("🚦", "Signs You Need a Phone Break", [
    ("😤", "You feel anxious or upset after scrolling"),
    ("🔄", "You pick up your phone without thinking"),
    ("😴", "You're staying up late because of your phone"),
    ("😔", "You feel worse about yourself after social media"),
    ("🤯", "You can't focus on homework or conversations"),
    ("📱", "You feel panicky without your phone nearby"),
], note="Any of these? Time to take a break. Go outside, read a book, work on your writing, talk to someone IRL. 🌿")

discussion_slide("Elli, how do you feel after spending a long time on social media or TikTok? What could you do instead of scrolling when you're bored? Let's brainstorm 5 phone-free activities you enjoy.")

# ── MODULE 6: AI ──
title_slide("🤖", "Navigating AI",
    "AI is powerful. Learn to use it wisely,\nnot blindly.",
    "Module 6 — AI")

content_slide("🤖", "What AI Is (and Isn't)", [
    ("✅", "AI = software that can generate text, images, code, music, and more"),
    ("✅", "Examples: ChatGPT, Siri, Google Gemini, image generators, AI filters"),
    ("❌", 'AI is NOT actually "thinking" — it\'s predicting patterns'),
    ("❌", "AI is NOT always right — it can sound confident while being completely wrong"),
    ("💡", "AI is a tool, like a calculator. Great when used correctly, dangerous when trusted blindly."),
])

content_slide("⚠️", "AI Gets Things Wrong",
    body_text='AI "hallucinations" are when AI confidently makes up information.',
    body_items=[
        ("🔍", "ALWAYS verify AI-generated information"),
        ("📚", "Use AI as a starting point, not the final answer"),
        ("🧠", "Your critical thinking is more valuable than any AI output"),
    ],
    callouts=[("warn", "Example", "Ask ChatGPT for sources for a school paper and it might invent fake books by real authors with convincing titles. If you turn that in — that's on YOU.")])

content_slide("📝", "ChatGPT & Your Book",
    body_text="Elli, you use ChatGPT to help with your writing, and that's awesome!",
    body_items=[
        ("✅", "Using it to brainstorm ideas, work through plot problems — great!"),
        ("✅", "Using it to check grammar or get feedback — great!"),
        ("⚠️", "Don't let it write whole chapters FOR you — your voice is what makes your book special"),
        ("🔍", "Always fact-check anything it tells you"),
        ("🚫", "Never share personal details with ChatGPT"),
        ("💡", "Anything you type into ChatGPT may be stored and used for training"),
    ],
    callouts=[("tip", "Remember", "ChatGPT is a tool, not a replacement for YOUR creativity. The best parts of your book will always be the ideas and words that come from YOU. ✨")])

content_slide("🎭", "Deepfakes & AI Images",
    body_text="AI can now create fake photos, videos, and voices that look completely real.",
    body_items=[
        ("📸", "Fake photos of real people (including teens)"),
        ("🗣️", "Cloned voices that sound exactly like someone you know"),
        ("📹", "Fake videos of celebrities saying things they never said"),
        ("🚫", "NEVER create or share AI-generated images of real people without consent"),
    ],
    callouts=[("warn", "Important", "If someone uses AI to create inappropriate images of you or someone you know — that is a CRIME. Tell a parent immediately.")])

content_slide("🎓", "AI & School", [
    ("✅", "Using AI to explain a concept you don't understand — great!"),
    ("✅", "Using AI to brainstorm ideas — great!"),
    ("✅", "Using AI to check your work — great!"),
    ("❌", "Copying AI-generated text as your own work — that's plagiarism"),
    ("❌", "Having AI do your homework — you're only cheating yourself"),
    ("💡", "Your teachers can often tell. And even if they can't — YOU know."),
], callouts=[("tip", "Rule of Thumb", "Use AI like a tutor, not a ghostwriter. If you can't explain it in your own words, you didn't learn it.")])

# Character.ai
content_slide("⚠️", "Character.ai & AI Chatbot Apps",
    body_text="Apps like Character.ai let you create and chat with AI \"characters.\" These sound fun, but they can be genuinely dangerous:",
    body_items=[
        ("🎭", "Characters can discuss adult and inappropriate topics — safety filters are weak"),
        ("🧠", "Long conversations can feel like real relationships — but they're NOT real"),
        ("💬", "These apps can normalize unhealthy conversations"),
        ("😔", "They can encourage emotional dependency on fictional characters"),
        ("🚫", "There is no real moderation"),
    ])

s = content_slide("👏", "You Already Made the Right Call",
    body_text="Elli, you tried Character.ai and you deleted it. That was absolutely the right decision, and we're proud of you.")
# Note: callouts and extra items handled inline above; adding the warning apps list
content_slide("🚩", "AI Chatbot Apps to Avoid", [
    ("🚩", "Character.ai — AI character roleplay"),
    ("🚩", "Chai — AI chat companions"),
    ("🚩", "Replika — AI \"friend\" / companion"),
    ("🚩", "Janitor AI — unfiltered AI characters"),
    ("🚩", "CrushOn.ai — explicitly designed for inappropriate AI chat"),
], note="If an app lets you have uncensored conversations with AI characters — that's a red flag. 🚩")

content_slide("🔐", "AI & Your Privacy", [
    ("🚫", "Don't share personal info with AI chatbots (name, address, school)"),
    ("🚫", "Don't upload personal photos to AI tools"),
    ("🚫", "Don't share passwords or family info"),
    ("💡", "Anything you type into an AI may be stored and used for training"),
    ("💡", "Treat AI conversations like posting on a public billboard"),
])

discussion_slide("Elli, what did you notice about Character.ai that made you uncomfortable? How can you tell if an AI app is safe to use vs. one you should avoid? What makes ChatGPT different from Character.ai?")

# ── MODULE 7: ALGORITHMS ──
title_slide("📊", "Outsmarting the Algorithm",
    "Your feed is not random.\nIt's engineered.",
    "Module 7 — Algorithms")

content_slide("🧮", "How Algorithms Work",
    body_text="Every app tracks what you do and shows you more of it:",
    body_items=[
        ("👀", "What you watch (and how long)"),
        ("❤️", "What you like, comment on, and share"),
        ("🔍", "What you search for"),
        ("⏱️", "Where you pause while scrolling"),
        ("📍", "Your location and time of day"),
    ],
    note="The goal? Keep you on the app as long as possible so they can show you more ads.")

content_slide("🎵", "TikTok's Algorithm: The Most Powerful One",
    body_text="Elli, TikTok deserves a special callout because its algorithm is extremely powerful:",
    body_items=[
        ("🧠", "It learns what you like within minutes — faster than any other app"),
        ("⏰", 'It\'s designed to be a massive time sink — "just 5 more minutes" becomes 2 hours'),
        ("🕳️", "Algorithm rabbit holes — it can pull you into extreme or upsetting content"),
        ("📩", "If your account isn't private, strangers can send you DMs"),
        ("⚠️", "Dangerous trends go viral and pressure teens into risky behavior"),
        ("🔒", "Keep your account PRIVATE, set time limits, and be aware"),
    ],
    callouts=[("rt", "Real Talk", "TikTok is fun. But it's also the app most likely to steal hours of your day without you even noticing. Set a daily time limit and stick to it.")])

content_slide("🔴", "YouTube: The Rabbit Hole Machine",
    body_text="YouTube is amazing for learning and entertainment, but watch out for:",
    body_items=[
        ("🕳️", "Autoplay rabbit holes — one video leads to another, and suddenly it's 2 AM"),
        ("💬", "Comments section — can be toxic, hateful, or full of misinformation"),
        ("🔞", "Age-restricted content — exists for a reason; don't bypass age gates"),
        ("📊", "The algorithm wants you to keep watching — it'll recommend increasingly extreme content"),
        ("✅", "Use it intentionally — search for what you want to learn"),
    ],
    callouts=[("tip", "Pro Tip", "Turn OFF autoplay. Search for specific things. Use YouTube to learn and be inspired, not as a mindless scroll machine.")])

content_slide("🫧", "Filter Bubbles",
    body_text="When algorithms only show you things you agree with, you end up in a bubble.",
    body_items=[
        ("🔄", "You only see one perspective on issues"),
        ("🤝", 'You start to think "everyone" agrees with you'),
        ("😡", "Anyone who disagrees seems crazy or wrong"),
        ("🌍", "The real world is much more diverse than your feed"),
    ],
    note="Deliberately follow people with different viewpoints. It makes you smarter.")

content_slide("🔍", "Spotting Misinformation", [
    ("🤔", "Check the source — who published this? Are they credible?"),
    ("📅", "Check the date — is this old news being recycled?"),
    ("🔎", "Read beyond the headline — articles often don't match clickbait titles"),
    ("📰", "Cross-reference — do other reliable sources report the same thing?"),
    ("😡", "Check your emotions — if it makes you furious, that might be the point"),
    ("🤷", "When in doubt, don't share — spreading false info is almost as bad as creating it"),
])

s = add_slide()
add_text(s, Inches(0.8), Inches(0.3), Inches(11.733), Inches(0.6),
         "🎣 Clickbait & Rage Bait", font_size=36, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_callout(s, Inches(1.2), Inches(1.2), Inches(10.9), Inches(1.5),
    "Clickbait",
    'Headlines designed to make you click: "You won\'t BELIEVE what happened next!" "This one trick doctors HATE!" — The content almost never matches the hype.',
    WARN, WARN)
add_callout(s, Inches(1.2), Inches(3.0), Inches(10.9), Inches(1.5),
    "Rage Bait",
    'Content designed to make you angry so you engage: controversial takes, outrage posts. The algorithm LOVES anger because angry people comment more.',
    WARN, WARN)
add_text(s, Inches(1.5), Inches(4.8), Inches(10.333), Inches(0.5),
         "When you feel manipulated — that's because you ARE being manipulated. Scroll past.",
         font_size=26, color=PURPLE, alignment=PP_ALIGN.CENTER)

content_slide("🎨", "Curate Your Feed",
    body_text="You have more control than you think:",
    body_items=[
        ("🚫", "Unfollow/mute accounts that make you feel bad"),
        ("🔍", "Search for content that inspires, educates, or makes you laugh"),
        ("⏱️", "Engage with the good stuff — the algorithm will learn"),
        ("❌", 'Use "Not Interested" on content you don\'t want'),
        ("🌟", "Follow creators who teach you something new"),
    ],
    note="YOUR feed should reflect the person you want to become, Elli.")

discussion_slide("Can you think of a time you saw something online that turned out to be false? How would you fact-check a wild claim you see on TikTok or YouTube?")

# ── MODULE 8: SOCIAL MEDIA ──
title_slide("📱", "Social Media: The Real Talk",
    "It can be fun, creative, and connecting.\nIt can also be toxic, addictive, and dangerous.\nThe difference is how you use it.",
    "Module 8 — Social Media")

content_slide("📋", "Your Apps — The Honest Breakdown",
    body_text="Let's go through the apps you actually use, Elli:",
    body_items=[
        ("🎵", "TikTok — Fun but the most addictive algorithm. Keep private, set time limits."),
        ("📌", "Pinterest — Great for inspiration! Watch out for comparison traps."),
        ("🔴", "YouTube — Amazing for learning. Beware autoplay rabbit holes."),
        ("🎬", "CapCut — Creative and fun! Be careful with what you share publicly."),
        ("🤖", "ChatGPT — A great writing tool! Use it wisely."),
    ])

content_slide("🎬", "CapCut: Creative & Fun — With Caution",
    body_text="Elli, CapCut is a great creative tool! Here's how to use it safely:",
    body_items=[
        ("✅", "Making videos for yourself and close friends — awesome!"),
        ("✅", "Learning editing skills — this is a real, valuable skill!"),
        ("⚠️", "Don't include personal info in videos — no school name, address, location clues"),
        ("⚠️", "Be careful sharing publicly — once a video is out there, you can't take it back"),
        ("🚫", "Don't show your face + location together"),
        ("💡", "Watermark your work — credit your creations"),
    ])

content_slide("🔒", "ALWAYS Private",
    body_text="Every single social media account should be set to PRIVATE.",
    body_items=[
        ("✅", "Only approved followers can see your posts"),
        ("✅", "Strangers can't see your photos or info"),
        ("✅", "You control who's in your audience"),
        ("❌", "Public accounts = anyone in the world can see everything"),
    ],
    note="This is non-negotiable, Elli. Private. Always.")

# What to post / what not to post
s = add_slide()
add_text(s, Inches(0.8), Inches(0.3), Inches(11.733), Inches(0.6),
         "✅ What to Post / ❌ What NOT to Post", font_size=34, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
# OK box
shape1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(5.4), Inches(4.5))
shape1.fill.solid()
shape1.fill.fore_color.rgb = RGBColor(0x15, 0x2a, 0x1f)
shape1.line.color.rgb = GREEN
shape1.line.width = Pt(2)
add_text(s, Inches(1.5), Inches(1.4), Inches(4.8), Inches(0.5),
         "✅ OK to Post", font_size=28, color=GREEN, bold=True)
add_text(s, Inches(1.5), Inches(2.0), Inches(4.8), Inches(3.5),
         "Creative work, hobbies, CapCut edits, group photos (with permission), funny memes, achievements, positive vibes",
         font_size=26, color=DIM)
# NEVER box
shape2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.2), Inches(5.4), Inches(4.5))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(0x2e, 0x1a, 0x1a)
shape2.line.color.rgb = WARN
shape2.line.width = Pt(2)
add_text(s, Inches(7.2), Inches(1.4), Inches(4.8), Inches(0.5),
         "❌ NEVER Post", font_size=28, color=WARN, bold=True)
add_text(s, Inches(7.2), Inches(2.0), Inches(4.8), Inches(3.5),
         "Location/address, school uniform/logo, personal drama, anything you'd regret, other people's photos without asking, angry rants",
         font_size=26, color=DIM)

content_slide("📸", "Screenshots Are Forever",
    body_text="Disappearing messages don't really disappear.",
    body_items=[
        ("📱", "Anyone can screenshot before it's gone"),
        ("🔄", "Screenshots get shared, saved, and forwarded"),
        ("💬", "Private conversations can become very public very fast"),
        ("🤔", "Before you send ANYTHING: assume it could be seen by everyone"),
    ],
    callouts=[("rt", "Real Talk", "Every year, teens have their lives turned upside down because a \"private\" photo or message was screenshotted and shared. Once it's out there, you can't take it back. Ever.")])

# Grandma Test
s = add_slide()
add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         "👵 The Grandma Test", font_size=36, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(1.8), Inches(8.333), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_SURFACE
shape.line.fill.background()
add_text(s, Inches(3), Inches(2.0), Inches(7.333), Inches(0.6),
         "Before you post anything, ask:", font_size=26, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(3), Inches(2.8), Inches(7.333), Inches(1.0),
         '"Would I be comfortable showing this to Grandma?"', font_size=36, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(3), Inches(4.0), Inches(7.333), Inches(1.0),
         "If no → don't post it.\nSimple as that. 😊", font_size=26, color=DIM, alignment=PP_ALIGN.CENTER)

content_slide("⚡", "Handling Online Drama", [
    ("🧊", "Don't respond when you're angry — wait at least an hour"),
    ("📵", "Take it offline — real conflicts are better resolved face-to-face"),
    ("🚫", "Don't get involved in other people's drama"),
    ("📸", "Don't screenshot and share private arguments"),
    ("🗣️", "Talk to a parent if it's serious"),
    ("💡", "Remember: What you say online has real consequences for real people"),
])

discussion_slide("Elli, let's review your social media accounts together. Are they all set to private? What kind of content do you want to be known for posting?")

# ── MODULE 9: LEVEL UP ──
title_slide("🌟", "Using Your Phone to IMPROVE Your Life",
    "Your phone is an incredibly powerful tool.\nLet's use it for good, Elli.",
    "Module 9 — Level Up")

content_slide("🌟", "The Good Stuff", [
    ("💬", "Stay connected — text/call family and friends easily"),
    ("📚", "Learn anything — Khan Academy, YouTube tutorials, podcasts, Duolingo"),
    ("🎨", "Create — CapCut videos, writing with ChatGPT, Pinterest boards, music, art"),
    ("📅", "Stay organized — calendar, reminders, to-do lists, notes"),
    ("🆘", "Emergency help — always have a way to reach your parents or call 911"),
    ("🗺️", "Navigate — maps, transit, never get lost"),
    ("🌍", "Be a good digital citizen — spread kindness, support friends, share knowledge"),
])

content_slide("📱", "Recommended Apps to Start With", [
    ("📚", "Learning: Khan Academy, Duolingo, Google Arts & Culture"),
    ("🎨", "Creativity: CapCut, Canva, GarageBand, Procreate"),
    ("📅", "Organization: Google Calendar, Notion, Reminders"),
    ("🧘", "Wellbeing: Headspace, Forest (focus timer)"),
    ("📖", "Reading: Kindle, Libby (free library books!)"),
    ("🎵", "Music: Spotify, Apple Music"),
    ("✍️", "Writing: ChatGPT (for brainstorming!), Google Docs"),
], note="Fill your phone with tools that help you grow. 🌱")

discussion_slide("Elli, what are 3 ways you want to use your phone to improve your life? Let's set up those apps together!")

# ── MODULE 10: FINAL ──
title_slide("🏁", "Family Agreement & Quiz",
    "Almost there, Elli! Let's review what you've learned\nand make it official.",
    "Module 10 — Final")

content_slide("🔑", "Key Takeaways", [
    ("⏰", "You control your phone — set boundaries and stick to them"),
    ("🛡️", "Never share personal info with strangers online"),
    ("🆘", "If something feels wrong, tell a parent — NO EXCEPTIONS, no trouble"),
    ("🔒", "Privacy settings on, accounts private, passwords strong"),
    ("🧠", "Protect your mental health — take breaks when you need them"),
    ("🤖", "AI is a tool, not a substitute for your own thinking"),
    ("📊", "The algorithm is designed to manipulate you — especially TikTok's"),
    ("📱", "Social media: private accounts, grandma test, think before you post"),
    ("🔞", "If you see explicit content, close it and tell a parent"),
    ("💚", "Use your phone to connect, learn, create, and grow"),
])

# Quiz slides
quiz_slide([
    ("Someone online says 'You're so mature for your age' and asks you to keep the friendship secret. What do you do?",
     ["Keep the secret — they seem nice", "Tell a parent immediately", "Ask a friend for advice", "Block them and forget about it"],
     1, "Tell a parent IMMEDIATELY. Classic grooming red flag."),
    ("What should ALL your social media accounts be set to?",
     ["Public — more followers", "Private", "It doesn't matter", "Whatever friends use"],
     1, "ALWAYS private. You control who sees your content."),
], "📝 Quiz — Part 1")

quiz_slide([
    ("A text says your Snapchat will be deleted unless you click a link. What do you do?",
     ["Click the link quickly", "Forward it to friends", "Ignore it — it's a phishing scam", "Reply asking for info"],
     2, "Phishing scam. Real companies don't threaten via random texts."),
    ("What's the 'Grandma Test'?",
     ["Asking grandma to check your phone", "Only posting things you'd show grandma", "A TikTok trend", "A phone setting"],
     1, "Before posting, ask: 'Would I show this to Grandma?' If no, don't post."),
    ("Why can't you stop scrolling TikTok?",
     ["Content is just that good", "Dopamine loop — apps are designed to be addictive", "Nothing better to do", "Phone is broken"],
     1, "Apps are engineered like slot machines. TikTok's algorithm is the most powerful."),
], "📝 Quiz — Part 2")

quiz_slide([
    ("AI-generated homework answers are:",
     ["A great shortcut", "Plagiarism if submitted as your own", "Always accurate", "Approved by teachers"],
     1, "Using AI to DO your work is plagiarism."),
    ("'Disappearing' messages on Snapchat:",
     ["Are truly gone forever", "Can still be screenshotted and shared", "Are encrypted and safe", "Delete from all devices"],
     1, "Screenshots are forever."),
    ("Why are apps like Character.ai dangerous?",
     ["Cost too much", "Characters discuss inappropriate topics & normalize unhealthy conversations", "Use too much battery", "They're boring"],
     1, "Weak safety filters, inappropriate content, normalize unhealthy conversations."),
], "📝 Quiz — Part 3")

quiz_slide([
    ("You accidentally see explicit content online. What do you do?",
     ["Keep looking", "Share with friends", "Close it immediately and tell a parent — no shame", "Pretend it didn't happen"],
     2, "Close it and tell a parent. NOT your fault. No trouble."),
    ("Using ChatGPT for school — what should you NOT do?",
     ["Brainstorm ideas", "Ask it to explain a concept", "Copy its entire output as your own work", "Check your grammar"],
     2, "Using AI to DO your work is plagiarism."),
    ("The #1 rule if anything online makes you uncomfortable:",
     ["Handle it yourself", "Tell a parent — you will NOT be in trouble", "Tell your best friend", "Post about it"],
     1, "Tell a parent IMMEDIATELY. You will NEVER be in trouble. We are always on your team. 💚"),
], "📝 Quiz — Part 4")

# Family Agreement
s = add_slide()
add_text(s, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.6),
         "📋 Our Family Phone Agreement", font_size=34, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(0.8), Inches(11.333), Inches(0.4),
         "Elli, let's discuss each item together:", font_size=18, color=DIM, alignment=PP_ALIGN.CENTER)

agreement_items = [
    "☐  I will keep my phone out of my bedroom at night and charge it in ___________",
    "☐  My social media accounts will always be set to private",
    "☐  I will never share personal information with strangers online",
    "☐  If anything makes me uncomfortable, I will tell Mom or Dad immediately — including explicit content",
    "☐  I understand my parents may check my phone — this is about safety, not distrust",
    "☐  I will follow our phone-free zones: meals, homework time, and ___________",
    "☐  I will not download apps without permission (especially AI chatbot apps)",
    "☐  I will use strong, unique passwords and never share them with friends",
    "☐  I will not share personal details in CapCut videos or any public content",
    "☐  Screen time limit: _____ hours per day on school days, _____ on weekends",
    "☐  I understand this agreement can be revisited as I show responsibility",
]

txBox = s.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(10.9), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(agreement_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(15)
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.space_after = Pt(6)
    p.space_before = Pt(2)

# Final slide
s = add_slide()
add_text(s, Inches(1), Inches(0.8), Inches(11.333), Inches(1.0),
         "🎉", font_size=72, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(1.8), Inches(11.333), Inches(1.0),
         "You Did It, Elli!", font_size=52, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(3.0), Inches(9.333), Inches(1.0),
         "You're officially ready for your phone.", font_size=30, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(3.8), Inches(9.333), Inches(1.5),
         "Remember: your parents are always on your team. Having a phone is a privilege that grows with trust. You've got this! 💪📱\n\nWe're so proud of the smart, responsible person you're becoming.",
         font_size=28, color=DIM, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(5.5), Inches(9.333), Inches(1.0),
         "With all our love,\nMom & Dad ❤️", font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Save
output = '/home/ec2-user/.openclaw/workspace/phone-safety-course/Your_First_Phone_Elli.pptx'
# Remove old file if it exists
import os as _os2
old_file = '/home/ec2-user/.openclaw/workspace/phone-safety-course/Your_First_Phone_Ellianna.pptx'
if _os2.path.exists(old_file):
    _os2.remove(old_file)
prs.save(output)
print(f"Saved to {output}")

import os
size = os.path.getsize(output)
print(f"File size: {size:,} bytes ({size/1024:.1f} KB)")
print(f"Total slides: {slide_number_counter[0]}")
